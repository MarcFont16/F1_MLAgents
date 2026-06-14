from pptx import Presentation
from pptx.util import Inches, Pt
import csv
import os

ROOT = os.path.dirname(os.path.dirname(__file__))
CSV_PATH = os.path.join(ROOT, 'results', 'final_evaluation_results.csv')
README_PATH = os.path.join(ROOT, 'README.md')
OUT_PPTX = os.path.join(ROOT, 'F1_MLAgents_presentation_detailed.pptx')

prs = Presentation()

# Title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Autonomous F1 Driving"
subtitle.text = "Spa-Francorchamps, Unity ML-Agents, training and evaluation"

# Agenda slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Agenda"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "1. Project motivation and architecture\n"
    "2. Training configuration\n"
    "3. Key code components\n"
    "4. Experiments and results\n"
    "5. Demo and next steps\n"
)

# Architecture slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Architecture & Data Flow"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Unity simulates the F1 vehicle and track\n"
    "- F1Agent uses ML-Agents PPO to decide actions\n"
    "- TelemetrySender sends speed/steering/friction via UDP\n"
    "- Node.js server forwards data to browser dashboard\n"
    "- EvaluationManager exports results to CSV\n"
)

# Training config slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Training Configuration"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- PPO trainer with 10M steps total\n"
    "- batch_size 2048, buffer_size 20480\n"
    "- lr 0.0003, epsilon 0.2, beta 0.005\n"
    "- hidden_units 256, num_layers 3\n"
    "- normalize observations, gamma 0.993, horizon 128\n"
)

# Reward shaping slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Reward Shaping"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- +1.0 per checkpoint to incentivize progress\n"
    "- -0.0002 per step to motivate finishing laps\n"
    "- Grass penalty: -0.002, Gravel penalty: -0.005\n"
    "- Crash penalty: -0.5 / -1.0 and end episode on wall hit\n"
    "- Jerk penalty reduces sudden steering changes\n"
)

# F1Agent slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "F1Agent.cs - Agent Behavior"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Observations: noisy current speed only\n"
    "- Actions: continuous acceleration and steering\n"
    "- Off-track detection uses raycast against floor tags\n"
    "- Movement uses Translate + SphereCast for collision handling\n"
    "- AddReward balances speed, safety and smoothness\n"
)

# Anti-wobble slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Anti-Wobble Steering"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Cubic steering gives fine control near center\n"
    "- Deadzone ignores micro-movements on straights\n"
    "- Lerp smooths steering transitions\n"
    "- Result: less oscillation at high speed\n"
)

# RaceManager slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "RaceManager.cs - Timing and Curriculum"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Manages lap timer, sector times and best splits\n"
    "- Uses 3 sector triggers to calcular tiempos parciales\n"
    "- ResetRaceOnCrash() restaura l’estat després d’un accident\n"
    "- Currículum: aument de velocitat després de voltes perfectes\n"
)

# Evaluation slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "EvaluationManager.cs"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Runs weather-based evaluation protocol\n"
    "- Adjusts track friction for each condition\n"
    "- Records successes, crashes and lap times\n"
    "- Writes results to CSV for analysis\n"
)

# Experiments slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Experiment History"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Baseline_Agent_01: no safety mechanisms, local minima\n"
    "- Experiment_02_SafetyWalls: walls added, but curriculum shock\n"
    "- Experiment_03_ProgressiveCurriculum: stable learning with gradual speed increase\n"
)

# Read README excerpt if exists
if os.path.exists(README_PATH):
    with open(README_PATH, 'r', encoding='utf-8') as f:
        lines = f.read().splitlines()
    excerpt = '\n'.join(lines[:15])
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Project Overview from README"
    body = slide.shapes.placeholders[1].text_frame
    body.text = excerpt

# Results slide (table)
if os.path.exists(CSV_PATH):
    with open(CSV_PATH, 'r', encoding='utf-8') as f:
        reader = csv.reader(f)
        rows = list(reader)
    if len(rows) > 0:
        headers = rows[0]
        data = rows[1:]
        cols = len(headers)
        rows_count = len(data) + 1
        slide = prs.slides.add_slide(prs.slide_layouts[5])
        slide.shapes.title.text = 'Final Evaluation Results'
        left = Inches(0.5)
        top = Inches(1.4)
        width = Inches(9)
        height = Inches(3.2)
        table = slide.shapes.add_table(rows_count, cols, left, top, width, height).table
        for c, h in enumerate(headers):
            table.cell(0, c).text = h
            para = table.cell(0, c).text_frame.paragraphs[0]
            para.font.bold = True
            para.font.size = Pt(12)
        for r, row in enumerate(data, start=1):
            for c, val in enumerate(row):
                table.cell(r, c).text = val

# Demo slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Demo Checklist"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Start Node.js server in F1_Dashboard\n"
    "- Open browser at http://localhost:3000\n"
    "- Run Unity scene with F1Agent and TelemetrySender\n"
    "- Show CSV results and explain success rates\n"
)

# Next steps slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
slide.shapes.title.text = "Next Steps"
body = slide.shapes.placeholders[1].text_frame
body.text = (
    "- Add more sensor observations (raycasts, IMU, distances)\n"
    "- Train longer and add more episodes\n"
    "- Improve domain randomization for track conditions\n"
    "- Explore sim-to-real transfer strategies\n"
)

prs.save(OUT_PPTX)
print(f'Generated {OUT_PPTX}')
